import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import tempfile

def main():
    if len(sys.argv) < 3:
        print("Usage: python pdf2pptx.py <input.pdf> <output.pptx>")
        sys.exit(1)

    pdf_file = sys.argv[1]
    pptx_file = sys.argv[2]

    try:
        # Create a new presentation
        prs = Presentation()
        # Set slide width and height to 4:3 (default) or we can dynamically adjust based on PDF size
        # We will use 16:9 which is more modern, or we can just use the default slide size.
        # Let's adjust slide size to match the first PDF page.
        
        pdf_document = fitz.open(pdf_file)
        
        if len(pdf_document) == 0:
            print("Error: Empty PDF")
            sys.exit(1)

        # Get dimensions of first page
        first_page = pdf_document[0]
        rect = first_page.rect
        width = rect.width
        height = rect.height
        
        # PyMuPDF returns points (72 points = 1 inch). pptx uses its own units.
        prs.slide_width = int((width / 72.0) * 914400)
        prs.slide_height = int((height / 72.0) * 914400)

        # Use a blank slide layout
        blank_slide_layout = prs.slide_layouts[6] 

        with tempfile.TemporaryDirectory() as temp_dir:
            for i in range(len(pdf_document)):
                page = pdf_document[i]
                # Render page to an image (dpi=150)
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                
                img_path = os.path.join(temp_dir, f"page_{i}.png")
                pix.save(img_path)
                
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add image taking up the full slide
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        pdf_document.close()
        prs.save(pptx_file)
        print("Success")

    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
